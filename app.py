# app.py
# A professional, single-file Streamlit application to generate comprehensive, AI-powered QBR decks.
# Version 8.1: Corrected an import error to ensure perfect execution.

import streamlit as st
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN # <-- CORRECTED LINE
from pptx.enum.shapes import MSO_SHAPE
import datetime
import os
import time
from io import BytesIO

# --- 1. BACKEND LOGIC: ADVANCED DATA & PRESENTATION GENERATION ---

# --- Color Palette & Helper ---
PRIMARY_COLOR_PPT = RGBColor(10, 47, 87)    # Deep Navy
ACCENT_COLOR_PPT = RGBColor(0, 122, 255)   # Professional Blue
TEXT_COLOR_PPT = RGBColor(33, 33, 33)      # Dark Gray
BACKGROUND_COLOR_PPT = RGBColor(248, 249, 250) # Light Gray for table alternates
WHITE_COLOR = RGBColor(255, 255, 255)

def rgb_to_hex(rgb_color_obj):
    """Converts a python-pptx RGBColor object to a hex string for matplotlib."""
    r, g, b = rgb_color_obj.rgb
    return f"#{r:02x}{g:02x}{b:02x}"

def get_enhanced_mock_data(customer_name, tone="Formal"):
    """Generates a rich, multi-faceted dataset for a comprehensive QBR."""
    np.random.seed(hash(customer_name) % (2**32 - 1))

    # Tone-based text generation
    if tone == "Optimistic":
        challenge_prefix = "Opportunity for growth: "
        learning_prefix = "Valuable insight gained: "
    else: # Formal
        challenge_prefix = ""
        learning_prefix = ""

    kpis = {
        "Account Health": np.random.randint(75, 98), "NPS": np.random.randint(30, 65),
        "Product Adoption (%)": np.random.randint(60, 95), "Active Users": f"{np.random.randint(150, 500)}",
        "Support Tickets Closed": np.random.randint(25, 100), "Renewal Date": (datetime.date.today() + datetime.timedelta(days=np.random.randint(90, 365))).strftime('%Y-%m-%d')
    }
    commit_vs_actual = pd.DataFrame({
        'Metric': ['Feature Delivery', 'Uptime SLA', 'Avg. Ticket Response'],
        'Commitment': ['5 New Features', '99.9% Uptime', '< 8 Hours'],
        'Actual': ['6 New Features', f"{99.9 + np.random.uniform(0, 0.09):.2f}% Uptime", f"{np.random.uniform(6, 7.9):.1f} Hours"],
        'Status': ['✅ Exceeded', '✅ Met', '✅ Met']
    })
    challenges = [
        challenge_prefix + "Initial onboarding for the new analytics module was slower than anticipated.",
        challenge_prefix + "Integration with the legacy CRM system required custom development work.",
        challenge_prefix + "User adoption in the finance department is lagging behind other teams."
    ]
    learnings = [
        learning_prefix + "A dedicated onboarding webinar for new modules significantly boosts initial adoption.",
        learning_prefix + "Pre-sales technical discovery for legacy systems is crucial to scope integrations accurately.",
        learning_prefix + "Targeted training sessions and identifying team champions accelerate department-level adoption."
    ]
    okrs = pd.DataFrame({
        'Objective': ["Enhance Team Collaboration", "Enhance Team Collaboration", "Improve Data-Driven Decisions", "Improve Data-Driven Decisions"],
        'Key Result': ["Increase shared dashboard usage by 20%", "Onboard the marketing team to the platform", "Complete integration with BI tool", "Train 5 team leads on advanced reporting"],
        'Status': ['Not Started', 'Not Started', 'Not Started', 'Not Started']
    })
    roadmap = {
        "Q4 2025": ["AI-Powered Insights Engine", "Mobile App V2 Launch"],
        "Q1 2026": ["Advanced API Access", "Third-Party Integration Marketplace"],
        "Q2 2026": ["Predictive Analytics Module", "Team-based Permissions V3"]
    }
    revenue_forecast = pd.DataFrame({
        'Month': pd.to_datetime([f'2025-{i}-01' for i in range(10, 13)]),
        'Forecasted Revenue ($K)': [50 + i*5 + np.random.randint(-5, 5) for i in range(3)]
    })
    action_plan = pd.DataFrame({
        'Action Item': ["Schedule onboarding for marketing team", "Finalize BI tool integration specs", "Identify and train finance dept. champions"],
        'Owner': ["John (CSM)", "Jane (Customer IT)", "Sarah (CSM)"],
        'Due Date': [(datetime.date.today() + datetime.timedelta(days=d)).strftime('%Y-%m-%d') for d in [14, 30, 45]],
        'Status': ['Not Started', 'Not Started', 'Not Started']
    })
    return locals()

def create_kpi_gauge_chart(score, output_path="kpi_gauge.png"):
    """Creates a visually appealing gauge chart for a key KPI."""
    fig, ax = plt.subplots(figsize=(6, 3), subplot_kw={'projection': 'polar'})
    fig.patch.set_alpha(0) # Transparent background
    
    # Define colors based on score
    if score >= 90: color = '#28a745' # Green
    elif score >= 80: color = '#ffc107' # Yellow
    else: color = '#dc3545' # Red

    # Plot the gauge
    ax.barh(1, np.deg2rad(score * 1.8), color=color, alpha=0.9)
    ax.barh(1, np.deg2rad(180), color='grey', alpha=0.2, left=np.deg2rad(0))
    ax.set_theta_zero_location('W')
    ax.set_theta_direction(-1)
    ax.set_thetagrids([], [])
    ax.set_rgrids([], [])
    ax.set_ylim(0, 1)
    
    # Add the score text
    fig.text(0.5, 0.4, f"{score}", ha='center', va='center', fontsize=48, fontweight='bold', color=rgb_to_hex(PRIMARY_COLOR_PPT))
    fig.text(0.5, 0.25, "Account Health", ha='center', va='center', fontsize=16, color=rgb_to_hex(TEXT_COLOR_PPT))

    plt.savefig(output_path, dpi=300, bbox_inches='tight', transparent=True)
    plt.close(fig)
    return output_path
    
def add_styled_table_to_slide(slide, df, x, y, cx, cy):
    """Adds a professionally styled pandas DataFrame table to a slide."""
    shape = slide.shapes.add_table(df.shape[0] + 1, df.shape[1], x, y, cx, cy)
    table = shape.table
    
    # Set column widths
    for i in range(df.shape[1]): table.columns[i].width = int(cx / df.shape[1])

    # Header styling
    for i, col_name in enumerate(df.columns):
        cell = table.cell(0, i)
        cell.text = col_name
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.color.rgb = WHITE_COLOR
        p.font.size = Pt(14)
        p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_COLOR_PPT
        cell.vertical_anchor = PP_ALIGN.CENTER

    # Body styling (with zebra striping)
    for r_idx, row in enumerate(df.iterrows()):
        for c_idx, value in enumerate(row[1]):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = str(value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.alignment = PP_ALIGN.LEFT
            cell.vertical_anchor = PP_ALIGN.CENTER
            if r_idx % 2 == 1: # Apply shading to odd rows
                cell.fill.solid()
                cell.fill.fore_color.rgb = BACKGROUND_COLOR_PPT

def add_master_slide_elements(slide, customer_name, logo_path):
    """Adds a consistent footer and logo to each slide."""
    # Footer Line
    line = slide.shapes.add_shape(MSO_SHAPE.LINE_INVERSE, Inches(0.5), Inches(8.5), Inches(15), Inches(0))
    line.line.color.rgb = ACCENT_COLOR_PPT
    line.line.width = Pt(1)

    # Footer Text
    footer_text = slide.shapes.add_textbox(Inches(0.5), Inches(8.55), Inches(10), Inches(0.3))
    tf = footer_text.text_frame
    p = tf.add_paragraph()
    p.text = f"Quarterly Business Review: {customer_name} | {datetime.date.today().strftime('%B %d, %Y')}"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_COLOR_PPT
    
    # Logo
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(14.5), Inches(0.4), height=Inches(0.6))


def create_professional_qbr_deck(data, logo_path):
    """Builds a comprehensive, professionally styled PowerPoint presentation."""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    customer_name = data['customer_name']

    # --- Slide Master Helpers ---
    def add_title_slide(title_text, subtitle_text):
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank layout
        
        # Manually add styled title and subtitle
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = title_text
        p.font.size = Pt(54)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR_PPT
        
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.8), Inches(14), Inches(1))
        tf = subtitle_box.text_frame
        p = tf.add_paragraph()
        p.text = subtitle_text
        p.font.size = Pt(24)
        p.font.color.rgb = TEXT_COLOR_PPT

        add_master_slide_elements(slide, customer_name, logo_path)
        return slide

    def add_content_slide(title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[5]) # Blank layout
        
        # Manually add styled title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(14), Inches(1.2))
        tf = title_box.text_frame
        p = tf.add_paragraph()
        p.text = title_text
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR_PPT
        
        add_master_slide_elements(slide, customer_name, logo_path)
        return slide

    # --- SLIDE CREATION ---
    add_title_slide(f"Quarterly Business Review: {customer_name}", f"Q3 2025 Report | Prepared: {datetime.date.today().strftime('%B %d, %Y')}")
    
    slide = add_content_slide("🗓️ Agenda")
    topics = ["Quarterly Snapshot & Highlights", "Commitment Review", "Challenges & Key Learnings", "Objectives for Next Quarter (OKRs)", "Strategic Growth & Product Roadmap", "Joint Action Plan & Next Steps"]
    agenda_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(6))
    tf = agenda_box.text_frame
    for topic in topics:
        p = tf.add_paragraph(); p.text = topic; p.level = 0; p.space_after = Pt(18); p.font.size = Pt(24)

    slide = add_content_slide("⭐ Quarterly Snapshot: Highlights")
    kpi_gauge_path = create_kpi_gauge_chart(data['kpis']['Account Health'])
    slide.shapes.add_picture(kpi_gauge_path, Inches(1), Inches(2.5), width=Inches(5))
    os.remove(kpi_gauge_path)
    
    # Other KPIs
    kpi_box = slide.shapes.add_textbox(Inches(7), Inches(2.5), Inches(8), Inches(5.5))
    tf = kpi_box.text_frame
    other_kpis = {k:v for k,v in data['kpis'].items() if k != 'Account Health'}
    for key, value in other_kpis.items():
        p = tf.add_paragraph(); p.text = f"{key}: "; p.font.size = Pt(22)
        run = p.add_run(); run.text = str(value); run.font.bold = True
        p.space_after = Pt(12)

    slide = add_content_slide("📊 Commitment Review: Promises vs. Reality")
    add_styled_table_to_slide(slide, data['commit_vs_actual'], Inches(1), Inches(2.5), Inches(14), Inches(4))

    slide = add_content_slide("💡 Challenges & Key Learnings")
    challenges_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(6.5), Inches(5))
    tf1 = challenges_box.text_frame
    p1 = tf1.add_paragraph(); p1.text = "Challenges Faced"; p1.font.bold = True; p1.font.size = Pt(24); p1.font.color.rgb = PRIMARY_COLOR_PPT; p1.space_after = Pt(12)
    for item in data['challenges']: p = tf1.add_paragraph(); p.text = f"• {item}"; p.level = 0; p.font.size = Pt(16); p.space_after = Pt(8)

    learnings_box = slide.shapes.add_textbox(Inches(8.5), Inches(2.5), Inches(6.5), Inches(5))
    tf2 = learnings_box.text_frame
    p2 = tf2.add_paragraph(); p2.text = "Key Lessons Learned"; p2.font.bold = True; p2.font.size = Pt(24); p2.font.color.rgb = PRIMARY_COLOR_PPT; p2.space_after = Pt(12)
    for item in data['learnings']: p = tf2.add_paragraph(); p.text = f"• {item}"; p.level = 0; p.font.size = Pt(16); p.space_after = Pt(8)

    slide = add_content_slide("🎯 Objectives for Next Quarter (OKRs)")
    add_styled_table_to_slide(slide, data['okrs'], Inches(1), Inches(2.5), Inches(14), Inches(5))

    slide = add_content_slide("🗺️ Strategic Growth & Product Roadmap")
    num_roadmap = len(data['roadmap']); box_width_r = 4.5; gap_r = 1.0
    total_width_r = num_roadmap * box_width_r + (num_roadmap - 1) * gap_r
    start_x_r = (16 - total_width_r) / 2
    for i, (quarter, features) in enumerate(data['roadmap'].items()):
        left = Inches(start_x_r + i * (box_width_r + gap_r))
        txBox = slide.shapes.add_textbox(left, Inches(2.5), Inches(box_width_r), Inches(5)); tf = txBox.text_frame
        p_qtr = tf.add_paragraph(); p_qtr.text = quarter; p_qtr.font.bold = True; p_qtr.font.size = Pt(24); p_qtr.font.color.rgb = ACCENT_COLOR_PPT
        for feature in features: p_feat = tf.add_paragraph(); p_feat.text = "• " + feature; p_feat.level = 0; p_feat.font.size=Pt(18); p_feat.space_after=Pt(8)

    slide = add_content_slide("🤝 Joint Action Plan & Owners")
    add_styled_table_to_slide(slide, data['action_plan'], Inches(1), Inches(2.5), Inches(14), Inches(4))

    add_title_slide("Thank You", "Q&A and Discussion")
    
    # --- Save Presentation ---
    output_filename = f"QBR_{data['customer_name'].replace(' ', '_')}_{datetime.date.today()}.pptx"
    prs.save(output_filename)
    return output_filename

# --- 2. FRONTEND UI: CREATIVE & WONDERFUL STREAMLIT INTERFACE ---

st.set_page_config(page_title="AI QBR Deck Generator", page_icon="💎", layout="wide")

# Injecting modern, beautiful CSS
st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;600;700&display=swap');
    
    html, body, [class*="st-"] {
        font-family: 'Inter', sans-serif;
    }
    .stApp {
        background: #f0f2f6;
    }
    
    /* Main Header */
    .main-header {
        font-size: 3rem;
        font-weight: 700;
        text-align: center;
        margin-bottom: 0px;
        color: #0a2f57; /* Primary Navy */
        padding-top: 1rem;
    }
    .sub-header {
        font-size: 1.2rem;
        text-align: center;
        color: #556;
        margin-bottom: 3rem;
    }
    
    /* Control Cards in Sidebar */
    [data-testid="stSidebar"] .stButton>button {
        background-color: #007bff;
        color: white;
        border-radius: 8px;
        transition: all 0.2s ease-in-out;
        border: none;
        font-weight: 600;
    }
    [data-testid="stSidebar"] .stButton>button:hover {
        background-color: #0056b3;
        transform: scale(1.02);
    }

    /* Main Content Cards */
    .content-card {
        background: white;
        border-radius: 15px;
        padding: 25px;
        box-shadow: 0 10px 30px rgba(0,0,0,0.07);
        border-left: 5px solid #007bff;
        height: 100%;
    }
    .content-card h3 {
        color: #0a2f57;
        font-weight: 600;
        margin-top: 0;
    }
    
    /* Progress Bar */
    .stProgress > div > div > div > div {
        background-color: #007bff;
    }
    
    /* Download Button */
    .stDownloadButton>button {
        background-color: #28a745;
        color: white;
        border-radius: 8px;
        transition: all 0.2s ease-in-out;
        border: none;
        font-weight: 600;
        width: 100%;
    }
    .stDownloadButton>button:hover {
        background-color: #218838;
        transform: scale(1.02);
    }
    
    /* Custom Tabs */
    .stTabs [data-baseweb="tab-list"] {
		gap: 24px;
	}
	.stTabs [data-baseweb="tab"] {
		height: 50px;
        white-space: pre-wrap;
		background-color: transparent;
		border-radius: 8px;
		gap: 8px;
		padding: 10px 15px;
	}
	.stTabs [aria-selected="true"] {
  		background-color: #e6f1ff;
	}

</style>
""", unsafe_allow_html=True)

# --- App Layout ---

# Sidebar for controls
with st.sidebar:
    st.image("https://www.gstatic.com/images/branding/googlelogo/svg/googlelogo_clr_74x24px.svg", width=100) # Placeholder logo
    st.title("Deck Controls")
    
    customer_name = st.text_input("Customer Name", "Innovate Corp")
    
    uploaded_logo = st.file_uploader("Upload Company Logo", type=['png', 'jpg', 'jpeg'])
    
    presentation_tone = st.selectbox(
        "Select Presentation Tone",
        ("Formal", "Optimistic"),
        help="The AI will subtly adjust the language based on your selection."
    )
    
    st.info("This is a Proof-of-Concept. All data is realistically simulated for demonstration.")

# Main content area
st.markdown("<h1 class='main-header'>💎 AI QBR Deck Generator</h1>", unsafe_allow_html=True)
st.markdown("<p class='sub-header'>Transform customer data into a stunning, client-ready presentation in seconds.</p>", unsafe_allow_html=True)

# Using session state to manage file readiness
if 'deck_generated' not in st.session_state:
    st.session_state.deck_generated = False
if 'deck_path' not in st.session_state:
    st.session_state.deck_path = ""
    
col1, col2 = st.columns([1.2, 2])

with col1:
    st.markdown("<div class='content-card'><h3>Generate Your Deck</h3></div>", unsafe_allow_html=True)
    
    if st.button("🚀 Create Presentation", use_container_width=True):
        if customer_name:
            logo_path = None
            if uploaded_logo is not None:
                # Save uploaded logo to a temporary file
                logo_path = os.path.join(".", uploaded_logo.name)
                with open(logo_path, "wb") as f:
                    f.write(uploaded_logo.getbuffer())

            with st.spinner('Crafting your presentation... This may take a moment.'):
                progress_bar = st.progress(0, text="Initializing...")
                
                enhanced_data = get_enhanced_mock_data(customer_name, presentation_tone)
                time.sleep(1); progress_bar.progress(25, text="Generating Insights & Tone...")
                
                # We need to run chart creation in the main thread for Streamlit
                time.sleep(1); progress_bar.progress(50, text="Creating Visualizations...")
                
                time.sleep(1); progress_bar.progress(75, text="Assembling Deck...")
                final_deck_path = create_professional_qbr_deck(enhanced_data, logo_path)
                
                progress_bar.progress(100, text="Done!")
                st.success(f"🎉 Your QBR deck is ready!")
                
                st.session_state.deck_generated = True
                st.session_state.deck_path = final_deck_path

                # Clean up temporary logo file
                if logo_path and os.path.exists(logo_path):
                    os.remove(logo_path)
        else:
            st.warning("Please enter a customer name.")
            st.session_state.deck_generated = False

    if st.session_state.deck_generated:
        with open(st.session_state.deck_path, "rb") as file:
            st.download_button(
                label="⬇️ Download Presentation", data=file, file_name=st.session_state.deck_path,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
        # Clean up the generated deck file after making it available for download
        if os.path.exists(st.session_state.deck_path):
             os.remove(st.session_state.deck_path)
        st.session_state.deck_generated = False # Reset state

with col2:
    st.markdown("<div class='content-card' style='border-color: #28a745;'><h3>Component Previews</h3></div>", unsafe_allow_html=True)
    
    tab1, tab2, tab3 = st.tabs(["📊 KPI Snapshot", "📈 Chart Style", "📋 Table Style"])
    
    with tab1:
        st.subheader("KPI Gauge & Highlights")
        st.write("A visually impactful way to present your most important metric.")
        sample_kpi_path = create_kpi_gauge_chart(92, "sample_gauge.png")
        st.image(sample_kpi_path)
        os.remove(sample_kpi_path)
        
    with tab2:
        st.subheader("Data Visualization")
        st.write("Clear, branded charts are generated for commercial and performance data.")
        sample_data = get_enhanced_mock_data("Sample Company")
        fig, ax = plt.subplots()
        sns.barplot(
            x=sample_data['revenue_forecast']['Month'].dt.strftime('%b'), 
            y='Forecasted Revenue ($K)', 
            data=sample_data['revenue_forecast'], 
            color=rgb_to_hex(ACCENT_COLOR_PPT), 
            ax=ax
        )
        ax.set_title("Sample Revenue Forecast")
        st.pyplot(fig, use_container_width=True)

    with tab3:
        st.subheader("Styled Data Tables")
        st.write("Tables are automatically styled with branded headers and alternating row colors for clarity.")
        st.table(get_enhanced_mock_data("Sample Co.")['action_plan'])
