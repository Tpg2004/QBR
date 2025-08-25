# qbr_generator.py
# The core logic for generating a QBR presentation.

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches
import datetime
import os

# --- 1. DATA SIMULATION (In a real app, this would be an API call to your CRM/database) ---
def get_mock_customer_data(customer_name):
    """Generates fake but realistic data for a fictional customer."""
    np.random.seed(42) # for consistent results
    
    # Monthly Active Users (MAU) - showing growth
    months = pd.to_datetime([f'2025-{i}-01' for i in range(5, 9)])
    mau_data = [150 + i*15 + np.random.randint(-10, 10) for i in range(4)]
    
    # Support Tickets
    tickets_data = {
        'Category': ['Technical', 'Billing', 'Feature Request', 'Technical'],
        'Status': ['Closed', 'Closed', 'Open', 'Closed'],
        'Resolution Time (Hours)': [4, 8, np.nan, 6]
    }
    
    # Key Performance Indicators (KPIs)
    kpis = {
        "Account Health Score": np.random.randint(85, 95),
        "Last Quarter NPS": np.random.randint(40, 60),
        "Product Adoption Rate (%)": 75 + np.random.randint(0, 10),
        "Renewal Date": "2026-03-01"
    }
    
    return {
        "customer_name": customer_name,
        "mau_df": pd.DataFrame({'Month': months, 'ActiveUsers': mau_data}),
        "tickets_df": pd.DataFrame(tickets_data),
        "kpis": kpis
    }

# --- 2. AI-POWERED INSIGHTS (Simulating Natural Language Generation) ---
def generate_ai_summary(data):
    """Generates a text summary from the data. A real app would use an LLM like GPT-4."""
    kpis = data['kpis']
    mau_df = data['mau_df']
    
    # Analyze MAU trend
    mau_growth = ((mau_df['ActiveUsers'].iloc[-1] - mau_df['ActiveUsers'].iloc[0]) / mau_df['ActiveUsers'].iloc[0]) * 100
    
    summary = (
        f"This quarter, {data['customer_name']} has demonstrated strong positive momentum. "
        f"The Account Health Score is a robust {kpis['Account Health Score']}/100, and product adoption remains high at {kpis['Product Adoption Rate (%)']}%. "
        f"We've observed a significant {mau_growth:.1f}% growth in Monthly Active Users over the period, indicating excellent engagement and value realization."
    )
    return summary

# --- 3. DATA VISUALIZATION ---
def create_mau_chart(mau_df, customer_name, output_path="mau_chart.png"):
    """Creates and saves a beautiful MAU trend chart."""
    plt.style.use('seaborn-v0_8-whitegrid')
    fig, ax = plt.subplots(figsize=(7, 4))
    
    sns.lineplot(x='Month', y='ActiveUsers', data=mau_df, marker='o', markersize=8, color='#4A90E2', ax=ax)
    
    # Aesthetic improvements
    ax.set_title(f'Monthly Active Users (MAU) Trend for {customer_name}', fontsize=14, weight='bold')
    ax.set_xlabel('Month', fontsize=10)
    ax.set_ylabel('Active Users', fontsize=10)
    ax.tick_params(axis='x', rotation=45)
    ax.grid(True, which='both', linestyle='--', linewidth=0.5)
    plt.tight_layout()
    
    # Save the chart
    plt.savefig(output_path, dpi=300)
    return output_path

# --- 4. PRESENTATION ASSEMBLY ---
def create_qbr_deck(data, summary, chart_path):
    """Builds the PowerPoint presentation from the data and assets."""
    prs = Presentation()
    
    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = f"Quarterly Business Review: {data['customer_name']}"
    subtitle.text = f"Q3 2025 Report\nPrepared on: {datetime.date.today().strftime('%B %d, %Y')}"

    # Slide 2: Executive Summary
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]
    title.text = "Executive Summary 🌟"
    content.text = summary

    # Slide 3: Key Performance Indicators
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Key Performance Indicators 📊"
    
    # Add KPIs in a two-column layout
    kpi_text_frame = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8.5), Inches(4)).text_frame
    kpi_text_frame.word_wrap = True
    for key, value in data['kpis'].items():
        p = kpi_text_frame.add_paragraph()
        p.text = f"{key}: "
        p.font.bold = True
        run = p.add_run()
        run.text = str(value)

    # Slide 4: User Engagement Trend
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "User Engagement Growth"
    slide.shapes.add_picture(chart_path, Inches(1.5), Inches(2.0), width=Inches(7))

    # Save the final presentation
    output_filename = f"QBR_{data['customer_name'].replace(' ', '_')}_{datetime.date.today()}.pptx"
    prs.save(output_filename)
    return output_filename

# --- Main function to orchestrate the process ---
def generate_qbr_for_customer(customer_name):
    """The main function that runs the entire QBR generation process."""
    print("1. Fetching customer data...")
    customer_data = get_mock_customer_data(customer_name)
    
    print("2. Generating AI-powered summary...")
    ai_summary = generate_ai_summary(customer_data)
    
    print("3. Creating data visualizations...")
    chart_file = create_mau_chart(customer_data['mau_df'], customer_name)
    
    print("4. Assembling the PowerPoint deck...")
    deck_filename = create_qbr_deck(customer_data, ai_summary, chart_file)
    
    print(f"\n✅ Success! Your QBR deck is ready: {deck_filename}")
    
    # Clean up the chart image
    os.remove(chart_file)
    
    return deck_filename

if __name__ == '__main__':
    # This allows you to run this script directly to test it
    generate_qbr_for_customer("Innovate Inc.")
